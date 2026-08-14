#!/usr/bin/env python3
"""
NEXTERP User Guide - PowerPoint Generator
Generates a professional PPTX presentation
Usage: python generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY_BLUE = RGBColor(37, 99, 235)
SECONDARY_GREEN = RGBColor(16, 185, 129)
TEXT_LIGHT = RGBColor(248, 250, 252)
TEXT_GRAY = RGBColor(148, 163, 184)
BG_DARK = RGBColor(15, 23, 42)


def add_dark_background(slide):
    """Add dark background to slide"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    return bg


def add_header(slide, title):
    """Add blue header bar with title"""
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT


def add_title_slide(title, subtitle):
    """Add a centered title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_GRAY
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(title, bullets):
    """Add slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, title)

    bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12), Inches(5.2))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(14)


def add_two_column_slide(title, left_items, right_items):
    """Add slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, title)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.8), Inches(5.2))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
            p.font.bold = True
            p.font.color.rgb = SECONDARY_GREEN
            p.font.size = Pt(20)
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(5.2))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
            p.font.bold = True
            p.font.color.rgb = PRIMARY_BLUE
            p.font.size = Pt(20)
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)


def add_table_slide(title, headers, rows):
    """Add slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, title)

    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(
        table_rows, cols, Inches(0.5), Inches(1.8),
        Inches(12.333), Inches(4.5)
    ).table

    # Header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_LIGHT


def add_step_slide(title, steps):
    """Add slide with numbered steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, title)

    for i, (step_title, step_desc) in enumerate(steps):
        y_pos = 1.7 + i * 1.2

        # Step number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos),
            Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = SECONDARY_GREEN
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.05), Inches(0.6), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER

        # Step title
        step_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos), Inches(11), Inches(0.5))
        tf = step_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT

        # Step description
        desc_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos + 0.45), Inches(11), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY


def add_workflow_slide(title, steps):
    """Add workflow diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_background(slide)
    add_header(slide, title)

    box_width = 2.2
    box_height = 1.0
    start_x = 0.8
    y = 3.8

    for i, step in enumerate(steps):
        x = start_x + i * (box_width + 0.8)

        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(box_width), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = PRIMARY_BLUE
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER

        # Arrow
        if i < len(steps) - 1:
            arrow_x = x + box_width + 0.1
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(y + 0.3),
                Inches(0.5), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = SECONDARY_GREEN
            arrow.line.fill.background()


# ============================================================================
# GENERATE PRESENTATION
# ============================================================================

# Slide 1: Cover
add_title_slide("NEXTERP", "User Guide - Enterprise Resource Planning System")

# Slide 2: Agenda
add_content_slide("Agenda", [
    "Getting Started - Login, Dashboard, Navigation",
    "Core Modules - Inventory, Sales, Accounting, HRM",
    "Additional Modules - Purchasing, Projects, Assets, Quality",
    "Reports & Analytics - Data insights and exports",
    "Best Practices - Tips for effective usage"
])

# Slide 3: Getting Started
add_step_slide("Getting Started", [
    ("Access NEXTERP", "Open browser and navigate to http://localhost:3000"),
    ("Login", "Enter username: admin, Password: Admin123!"),
    ("Dashboard", "View KPIs, charts, recent orders, and quick actions"),
    ("Navigation", "Use sidebar menu to access different modules")
])

# Slide 4: Dashboard Overview
add_content_slide("Dashboard Overview", [
    "KPI Cards - Revenue, Orders, Customers at a glance",
    "Charts - Visual trends and analytics",
    "Recent Activity - Latest transactions",
    "Quick Actions - Common operations",
    "Date Filtering - Today, This Week, This Month, Custom Range"
])

# Slide 5: Inventory Management
add_two_column_slide("Inventory Management",
    ["Warehouse Management", "Create multiple locations", "Set default warehouse", "Track stock levels"],
    ["Stock Items", "SKU, Name, Category", "Unit of Measure", "Standard Cost and Price", "Reorder Level"]
)

# Slide 6: Inventory Features
add_table_slide("Inventory Features",
    ["Feature", "Description"],
    [
        ["Warehouses", "Create and manage warehouse locations"],
        ["Stock Items", "Product catalog with SKU and pricing"],
        ["Stock Movements", "Track IN/OUT/Adjustment transactions"],
        ["Low Stock Alerts", "Automatic notifications"],
        ["Batch Tracking", "Track items by batch number"],
        ["Reports", "Stock valuation, movement history"]
    ]
)

# Slide 7: Sales Workflow
add_workflow_slide("Sales Order Workflow", [
    "DRAFT", "SUBMITTED", "APPROVED", "INVOICED", "PAID"
])

# Slide 8: Sales Management
add_two_column_slide("Sales Management",
    ["Customer Management", "Contact Information", "Credit Limits", "Payment Terms", "Order History"],
    ["Order Process", "1. Create Sales Order", "2. Submit for Approval", "3. Generate Invoice", "4. Record Payment"]
)

# Slide 9: Chart of Accounts
add_table_slide("Chart of Accounts",
    ["Type", "Nature", "Examples"],
    [
        ["Asset", "Debit increases", "Cash, Inventory, Equipment"],
        ["Liability", "Credit increases", "Accounts Payable, Loans"],
        ["Equity", "Credit increases", "Capital, Retained Earnings"],
        ["Revenue", "Credit increases", "Sales, Services"],
        ["Expense", "Debit increases", "Rent, Salary, Utilities"]
    ]
)

# Slide 10: Journal Entry
add_step_slide("Creating Journal Entry", [
    ("Select Entry Date", "Choose the transaction date"),
    ("Enter Reference", "Add reference number or description"),
    ("Add Line Items", "Select Account + Debit/Credit amount"),
    ("Verify Balance", "Ensure Debit = Credit"),
    ("Submit", "Submit for approval")
])

# Slide 11: Journal Workflow
add_workflow_slide("Journal Entry Workflow", [
    "DRAFT", "SUBMITTED", "POSTED"
])

# Slide 12: HRM
add_two_column_slide("HRM - Human Resources",
    ["Employee Management", "Personal Information", "Contact Details", "Department & Position", "Employment Type"],
    ["Attendance Tracking", "Check In / Check Out", "Work Hours Recording", "Status: Present, Absent, Late", "Leave Management"]
)

# Slide 13: Leave Management
add_step_slide("Leave Request Process", [
    ("Submit Request", "Employee submits leave request"),
    ("Manager Review", "Manager reviews the request"),
    ("Approve/Reject", "Decision is made"),
    ("Balance Update", "Leave balance is updated automatically")
])

# Slide 14: Additional Modules
add_content_slide("Additional Modules", [
    "Purchasing - Suppliers, Purchase Orders, Goods Receipt",
    "Projects - Project Planning, Task Tracking, Gantt Charts",
    "Assets - Fixed Assets, Depreciation, Maintenance",
    "Quality - Inspections, NCR, CAPA Process",
    "Analytics - Real-time Dashboards, KPI Tracking"
])

# Slide 15: Reports
add_table_slide("Reports & Analytics",
    ["Report Type", "Contents"],
    [
        ["Sales Reports", "Revenue, orders, customer analysis"],
        ["Inventory Reports", "Stock levels, movements, valuation"],
        ["Financial Reports", "Trial balance, P&L, Balance Sheet"],
        ["HR Reports", "Attendance, leave, headcount"]
    ]
)

# Slide 16: Export Options
add_content_slide("Export Options", [
    "CSV - For Excel/spreadsheet analysis",
    "PDF - For printing and sharing",
    "Print - Direct printing from browser"
])

# Slide 17: Best Practices
add_content_slide("Best Practices - Daily Operations", [
    "Check Dashboard - Monitor KPIs daily",
    "Review Alerts - Address issues promptly",
    "Verify Stock - Prevent stockouts",
    "Backup Data - Protect business data"
])

# Slide 18: Security
add_content_slide("Security Best Practices", [
    "Strong Passwords - Use complex, unique passwords",
    "Role-Based Access - Assign appropriate permissions",
    "Regular Updates - Keep system updated",
    "Audit Logs - Review regularly"
])

# Slide 19: Data Integrity
add_content_slide("Data Integrity", [
    "Always use approval workflows",
    "Never share login credentials",
    "Reconcile accounts monthly",
    "Document all corrections"
])

# Slide 20: Summary
add_content_slide("Summary - Key Takeaways", [
    "Single system for all business operations",
    "Real-time visibility across modules",
    "Automated workflows reduce manual work",
    "Role-based security protects data",
    "Built-in reports for decision making"
])

# Slide 21: Questions
add_title_slide("Questions?", "Contact: support@nexterp.com")

# Slide 22: Thank You
add_title_slide("Thank You!", "Start Using NEXTERP: http://localhost:3000")

# Save presentation
output_path = "docs/NEXTERP_USER_GUIDE.pptx"
prs.save(output_path)
print(f"SUCCESS: Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
